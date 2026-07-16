from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree
import copy, os

prs = Presentation("IA_na_Sala_de_Aula_CDC_2026.pptx")

DARK = "1A3A3A"
TERRACOTA = "C75B39"
SAGE = "5B8C5A"
WHITE = "FFFFFF"
GOLD = "D4A843"
LTEXT = "5A5A5A"

def h2r(h):
    return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))

def rshape(s, l, t, w, h, fill):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = h2r(fill); sh.line.fill.background()
    return sh

def rtext(s, l, t, w, h, items, **kw):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        txt, opts = item if isinstance(item,tuple) else (item,{})
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = txt
        p.font.name = opts.get("font",kw.get("font","Calibri"))
        p.font.size = opts.get("size",kw.get("size",Pt(14)))
        p.font.bold = opts.get("bold",kw.get("bold",False))
        p.font.italic = opts.get("italic",kw.get("italic",False))
        p.font.color.rgb = h2r(opts.get("color",kw.get("color",DARK)))
        p.alignment = opts.get("align",kw.get("align",PP_ALIGN.LEFT))
    return tb

# New slide using blank layout
blank_layout = prs.slide_layouts[0]  # DEFAULT layout
new_slide = prs.slides.add_slide(blank_layout)
# Remove existing placeholders
for ph in list(new_slide.shapes):
    sp = ph._element
    sp.getparent().remove(sp)

# Background
new_slide.background.fill.solid()
new_slide.background.fill.fore_color.rgb = h2r(DARK)

W, H = Inches(10), Inches(5.625)

# Gold bar
rshape(new_slide, Emu(0), Emu(0), W, Inches(0.06), GOLD)
# Decorative ovals
sh = new_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(3.5), Inches(4), Inches(4))
sh.fill.solid(); sh.fill.fore_color.rgb = h2r(TERRACOTA); sh.line.fill.background()
sh.fill.fore_color.brightness = 0.3
sh2 = new_slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1), Inches(3), Inches(3))
sh2.fill.solid(); sh2.fill.fore_color.rgb = h2r(SAGE); sh2.line.fill.background()
sh2.fill.fore_color.brightness = 0.4

# Title
rtext(new_slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7),
    [("Outras Ferramentas Deste Repositório", {"bold":True,"size":Pt(28),"color":WHITE,"font":"Georgia"})],
    align=PP_ALIGN.LEFT)

# Tool cards
tools = [
    ("docx/", "Cria e edita documentos Word (.docx) com formatação profissional"),
    ("xlsx/", "Cria e edita planilhas (.xlsx, .csv) com fórmulas, gráficos e formatação"),
    ("pdf/", "Extrai texto, combina, divide, aplica OCR e preenche formulários PDF"),
    ("pptx/", "Cria e edita apresentações de slides (.pptx) com design personalizado"),
    ("skill-creator/", "Cria, edita e testa skills com avaliação automatizada de desempenho"),
]

y0 = Inches(1.4)
for i, (nm, desc) in enumerate(tools):
    y = y0 + Inches(i * 0.75)
    rshape(new_slide, Inches(1.0), y, Inches(8.0), Inches(0.6), WHITE)
    rshape(new_slide, Inches(1.0), y, Inches(0.06), Inches(0.6), TERRACOTA)
    rtext(new_slide, Inches(1.2), y + Inches(0.05), Inches(2.0), Inches(0.5),
        [(nm, {"bold":True,"size":Pt(16),"color":DARK,"font":"Georgia"})])
    rtext(new_slide, Inches(3.0), y + Inches(0.05), Inches(5.8), Inches(0.5),
        [(desc, {"size":Pt(13),"color":LTEXT})])

# Footer
rtext(new_slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.4),
    [("Use a skill pptx/ para criar ou editar apresentações como esta!", 
      {"size":Pt(12),"color":GOLD,"italic":True})],
    align=PP_ALIGN.CENTER)

# Reorder: move new slide (last) to second-to-last
ns = prs.part._element.find('.//p:sldIdLst', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
children = list(ns)
# last child is our new slide
last = children[-1]
# remove it
ns.remove(last)
# insert before last
ns.insert(len(children)-2, last)

prs.save("IA_na_Sala_de_Aula_CDC_2026.pptx")
print("OK: slide inserido como penúltimo")
